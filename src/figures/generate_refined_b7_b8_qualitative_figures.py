"""Generate publication-ready 900-DPI qualitative B7-B8 figures.

Outputs:
- Figure 5: a clearly labelled strong held-out Muscovite-Granite-3 input/GT/prediction triplet.
- Figure 6: two quantified source-specific error rows (Muscovite-Schist-3 and Pyroxene-Peridotite-1).
- Editable PowerPoint source deck: text, labels, callouts, and legend are separate objects; microscopy,
  reference masks, and predictions are raster scientific panels.
"""
from __future__ import annotations
import argparse, csv, json
from pathlib import Path

import numpy as np
import torch
from PIL import Image
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from torchvision.models.segmentation import deeplabv3_resnet50
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from generate_b7_native_900dpi_figures import PretrainedResUNet, MEAN, STD, C, miou, rgb

CLASS_NAMES = ['background','olivine','pyroxene','plagioclase','alkali feldspar','quartz','biotite','muscovite','hornblende']
PALETTE = [(22,27,34),(83,170,112),(59,130,191),(104,167,207),(250,184,65),(192,126,77),(140,82,162),(210,108,164),(75,176,177)]
NAVY='#14506f'; TEAL='#168272'; RUST='#b9644e'; INK='#122534'; LIGHT='#f4f8f9'


def infer(model7, model8, image, weight, device):
    x=torch.from_numpy(np.asarray(image,dtype=np.float32).transpose(2,0,1)/255.).unsqueeze(0)
    x=(x-MEAN)/STD
    with torch.no_grad():
        return (weight*model7(x.to(device)).softmax(1)+(1-weight)*model8(x.to(device))['out'].softmax(1)).argmax(1).squeeze().cpu().numpy().astype(np.int64)


def load_records(manifest, b7, b8, weight, device):
    with Path(manifest).open(newline='',encoding='utf-8') as f:
        rows=[r for r in csv.DictReader(f) if r['split']=='test']
    records=[]
    for row in rows:
        with Image.open(row['image_path']) as im: image=im.convert('RGB')
        with Image.open(row['mask_path']) as ma: gt=np.asarray(ma.convert('L'),dtype=np.int64)
        pred=infer(b7,b8,image,weight,device)
        records.append({'source_id':row['source_id'],'image_path':row['image_path'],'image':image,'gt':gt,'pred':pred,'frame_miou':miou(gt,pred)})
    return records


def add_panel(ax, data, label, header, color):
    ax.imshow(data); ax.set_xticks([]); ax.set_yticks([])
    for spine in ax.spines.values(): spine.set_edgecolor(color); spine.set_linewidth(2.8)
    ax.add_patch(Rectangle((0,0),1,0.10,transform=ax.transAxes,facecolor=color,clip_on=False,zorder=3))
    ax.text(0.5,0.05,header,transform=ax.transAxes,ha='center',va='center',color='white',fontsize=14,fontweight='bold',zorder=4)
    ax.set_title(label,fontsize=14,pad=16,fontweight='bold',color=INK)


def save_figure5(record,out):
    fig,axs=plt.subplots(1,3,figsize=(15,6),constrained_layout=True)
    fig.patch.set_facecolor('white')
    fig.suptitle('Figure 5. Strong held-out B7–B8 ensemble illustration',fontsize=22,fontweight='bold',color=NAVY,y=1.05)
    fig.text(0.5,0.99,'Muscovite–Granite–3 is shown as a favourable held-out example; the primary performance claim remains grouped source-level cross-validation.',ha='center',fontsize=11,color='#475760')
    add_panel(axs[0],record['image'],'Muscovite–Granite–3, registered frame', 'POLARIZED-LIGHT INPUT',TEAL)
    add_panel(axs[1],rgb(record['gt']),'Reference semantic mask','GROUND TRUTH',RUST)
    add_panel(axs[2],rgb(record['pred']),'B7–B8 ensemble prediction','PREDICTION',RUST)
    fig.text(0.50,-0.02,f'Frame mineral mIoU = {record["frame_miou"]:.4f}  |  Native image resolution: 1280 × 1024  |  B7 = 0.6; B8 = 0.4',ha='center',fontsize=12,fontweight='bold',color=INK)
    legend='  '.join([f'{name}' for name in CLASS_NAMES[1:]])
    fig.text(0.50,-0.09,'Mineral classes: '+legend,ha='center',fontsize=9,color='#475760')
    png=out/'Figure_5_strong_native_ensemble_900dpi.png'; tif=out/'Figure_5_strong_native_ensemble_900dpi.tif'
    fig.savefig(png,dpi=900,bbox_inches='tight',facecolor='white'); fig.savefig(tif,dpi=900,bbox_inches='tight',facecolor='white'); plt.close(fig)
    return png,tif


def save_figure6(schist, peridotite, out):
    items=[
        (schist,'Muscovite–Schist–3','Mineral mIoU = 0.0866; composition L1 = 0.6848','Quartz/muscovite label fragmentation; reference quartz–muscovite interface: 1.000 → 0.155.'),
        (peridotite,'Pyroxene–Peridotite–1','Mineral mIoU = 0.0635; composition L1 = 0.5369','Dominant pyroxene → olivine reassignment; 71.7% of reference pyroxene predicted as olivine.'),
    ]
    fig,axs=plt.subplots(2,3,figsize=(15,12),constrained_layout=True)
    fig.patch.set_facecolor('white')
    fig.suptitle('Figure 6. Quantified source-specific B7–B8 error mechanisms',fontsize=22,fontweight='bold',color=NAVY,y=1.02)
    for row,(record,source,metric,mechanism) in enumerate(items):
        add_panel(axs[row,0],record['image'],source,'POLARIZED-LIGHT INPUT',TEAL)
        add_panel(axs[row,1],rgb(record['gt']),'Reference semantic mask','GROUND TRUTH',RUST)
        add_panel(axs[row,2],rgb(record['pred']),'B7–B8 ensemble prediction','PREDICTION',RUST)
        fig.text(0.50,0.49-row*0.47,metric,ha='center',fontsize=11,fontweight='bold',color=INK)
    fig.text(0.5,-0.02,'Each row is one held-out source; its eight registered angular views were aggregated for the source-level metrics. See Table 5 and Tables S2–S3.',ha='center',fontsize=10,color='#475760')
    png=out/'Figure_6_source_specific_errors_900dpi.png'; tif=out/'Figure_6_source_specific_errors_900dpi.tif'
    fig.savefig(png,dpi=900,bbox_inches='tight',facecolor='white'); fig.savefig(tif,dpi=900,bbox_inches='tight',facecolor='white'); plt.close(fig)
    return png,tif


def write_textbox(slide,left,top,width,height,text,size=16,bold=False,color=(18,37,52),align=PP_ALIGN.CENTER):
    shape=slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
    p=shape.text_frame.paragraphs[0]; p.alignment=align; run=p.add_run(); run.text=text; run.font.name='Arial'; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=RGBColor(*color)
    return shape


def add_case_slide(prs,title,subtitle,record_rows):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=RGBColor(255,255,255)
    write_textbox(slide,.35,.25,12.6,.45,title,25,True,(20,80,111),PP_ALIGN.LEFT)
    write_textbox(slide,.35,.72,12.6,.32,subtitle,11,False,(71,87,96),PP_ALIGN.LEFT)
    n=len(record_rows)
    # A single strong-case triplet may use large panels; two error rows require deliberately
    # compact panels and fixed vertical separation so editable metric callouts never overlap.
    if n == 1:
        geometry = [(1.20, .45, 4.25, 3.75, 3.75)]
    else:
        geometry = [(1.10, .45, 4.25, 2.00, 3.75), (4.20, .45, 4.25, 2.00, 3.75)]
    for r,(record,source,metric,mechanism) in enumerate(record_rows):
        y, first_x, step_x, picture_h, picture_w = geometry[r]
        labels=[('POLARIZED-LIGHT INPUT',record['image'],TEAL),('GROUND TRUTH',rgb(record['gt']),RUST),('PREDICTION',rgb(record['pred']),RUST)]
        for j,(label,img,color) in enumerate(labels):
            x=first_x+j*step_x
            img_path=Path('/tmp')/f'{source}_{r}_{j}.png'; img.save(img_path)
            slide.shapes.add_picture(str(img_path),Inches(x),Inches(y+.38),Inches(picture_w),Inches(picture_h))
            rect=slide.shapes.add_shape(1,Inches(x),Inches(y),Inches(picture_w),Inches(.32)); rect.fill.solid(); rect.fill.fore_color.rgb=RGBColor(*(tuple(int(color[i:i+2],16) for i in (1,3,5)))); rect.line.fill.background()
            write_textbox(slide,x,y+.02,picture_w,.22,label,9,True,(255,255,255))
        caption_y = y + .38 + picture_h + .05
        write_textbox(slide,.45,caption_y,12.1,.25,source+' — '+metric,11,True,(18,37,52))
        # Detailed transition mechanisms are retained in Table S2 and the figure caption;
        # omitting them here preserves a clean non-overlapping two-row visual layout.
    return slide


def build_pptx(figure5_rec, schist, peridotite,out):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    add_case_slide(prs,'Figure 5. Strong held-out B7–B8 ensemble illustration','Editable captions and labels; raster panels preserve native microscopy and semantic-map data.',[(figure5_rec,'Muscovite–Granite–3',f'Frame mineral mIoU = {figure5_rec["frame_miou"]:.4f}','Favourable illustration only; primary evidence is grouped source-level cross-validation.')])
    add_case_slide(prs,'Figure 6. Quantified source-specific error mechanisms','Each row aggregates eight registered angular views for the reported source-level metrics.',[(schist,'Muscovite–Schist–3','Mineral mIoU = 0.0866; composition L1 = 0.6848','Quartz/muscovite label fragmentation; interface 1.000 → 0.155.'),(peridotite,'Pyroxene–Peridotite–1','Mineral mIoU = 0.0635; composition L1 = 0.5369','71.7% of reference pyroxene reassigned to olivine.')])
    pptx=out/'Figures_5_6_Qualitative_Ensemble_EDITABLE.pptx'; prs.save(pptx); return pptx


def main():
    p=argparse.ArgumentParser(); p.add_argument('--manifest',required=True);p.add_argument('--b7-checkpoint',required=True);p.add_argument('--b8-checkpoint',required=True);p.add_argument('--b7-weight',type=float,default=.6);p.add_argument('--output',required=True);a=p.parse_args()
    if not torch.cuda.is_available(): raise RuntimeError('CUDA is required for native B7–B8 ensemble inference.')
    out=Path(a.output);out.mkdir(parents=True,exist_ok=True);dev=torch.device('cuda')
    b7=PretrainedResUNet().to(dev);b8=deeplabv3_resnet50(weights=None,weights_backbone=None,num_classes=C,aux_loss=False).to(dev)
    b7.load_state_dict(torch.load(a.b7_checkpoint,map_location=dev,weights_only=False)['model']);b8.load_state_dict(torch.load(a.b8_checkpoint,map_location=dev,weights_only=False)['model']);b7.eval();b8.eval()
    records=load_records(a.manifest,b7,b8,a.b7_weight,dev)
    def by_source(name):
        r=[x for x in records if x['source_id']==name]
        if not r: raise RuntimeError(f'Missing held-out source: {name}')
        return sorted(r,key=lambda z:z['frame_miou'],reverse=True)[0]
    strong=by_source('Muscovite-Granite-3');schist=by_source('Muscovite-Schist-3');peridotite=by_source('Pyroxene-Peridotite-1')
    f5=save_figure5(strong,out);f6=save_figure6(schist,peridotite,out);pptx=build_pptx(strong,schist,peridotite,out)
    summary={'dpi':900,'model':'validation-calibrated B7–B8 probability ensemble','b7_weight':a.b7_weight,'strong_frame':{'source':strong['source_id'],'image':strong['image_path'],'frame_miou':strong['frame_miou']},'source_error_figure_sources':['Muscovite-Schist-3','Pyroxene-Peridotite-1'],'figure_5':[str(x) for x in f5],'figure_6':[str(x) for x in f6],'editable_pptx':str(pptx)}
    (out/'refined_qualitative_figure_summary.json').write_text(json.dumps(summary,indent=2),encoding='utf-8');print(json.dumps(summary,indent=2))

if __name__=='__main__':main()
