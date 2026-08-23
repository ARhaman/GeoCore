"""Build original editable context figures for the mineral-segmentation manuscript.

This script produces four distinct PowerPoint figures:
  1. Dataset construction and source-level study design using real registered images and masks.
  2. Representative mineral-image and mask gallery selected from the real dataset.
  3. B7-B8 ensemble framework and validation-only model-selection logic.
  4. Development and grouped-validation evidence using audited study metrics.

All typography, boxes, arrows, equations, legends, and charts are editable PowerPoint objects.
Only real microscopy images and indexed mineral masks are raster panels. The figures are original
and do not reproduce external manuscript artwork.
"""
from __future__ import annotations

import argparse
import csv
import re
from collections import defaultdict
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAMES = ["background", "olivine", "pyroxene", "plagioclase", "alkali feldspar", "quartz", "biotite", "muscovite", "hornblende"]
PALETTE = [(32,35,40),(80,170,90),(70,125,200),(245,180,65),(190,120,70),(220,220,220),(140,80,155),(225,120,175),(70,175,175)]
NAVY=(18,73,110); TEAL=(26,126,111); BLUE=(50,120,168); RUST=(180,92,72); INK=(20,34,45); PALE=(244,248,250); MID=(75,91,102)


def frame_number(path: str) -> int:
    match = re.search(r"_(\d+)\.(?:jpg|jpeg|png)$", path.lower())
    return int(match.group(1)) if match else -1


def box(slide, x, y, w, h, fill, line=None, rounded=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*(line or fill))
    return shape


def text(slide, value, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    run = paragraph.add_run(); run.text = value; run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = RGBColor(*color)
    return shape


def arrow(slide, x1, y1, x2, y2, colour=TEAL, width=2.3):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(*colour); connector.line.width = Pt(width); connector.line.end_arrowhead = True
    return connector


def header(slide, number, title, subtitle):
    box(slide, 0, 0, 13.333, .55, NAVY, rounded=False)
    text(slide, f"Figure {number}. {title}", .27, .065, 12.6, .22, 20, (255,255,255), True)
    text(slide, subtitle, .30, .61, 12.2, .16, 9.5, MID)


def panel(slide, image_path, label, x, y, w, h, colour, label_size=9.0):
    box(slide, x, y, w, .24, colour, rounded=False)
    text(slide, label, x+.02, y+.04, w-.04, .13, label_size, (255,255,255), True, PP_ALIGN.CENTER)
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y+.28), width=Inches(w), height=Inches(h))


def rgb_mask(mask_path: Path, out_path: Path):
    with Image.open(mask_path) as source:
        values = np.asarray(source.convert("L"), dtype=np.uint8)
    result = np.zeros((*values.shape, 3), dtype=np.uint8)
    for index, colour in enumerate(PALETTE):
        result[values == index] = colour
    Image.fromarray(result, mode="RGB").save(out_path)


def source_records(rows):
    by_source = defaultdict(list)
    for row in rows:
        by_source[row["source_id"]].append(row)
    for values in by_source.values():
        values.sort(key=lambda row: frame_number(row["image_path"]))
    return by_source


def classes_present(mask_path: Path):
    with Image.open(mask_path) as source:
        values = np.asarray(source.convert("L"), dtype=np.uint8)
    return set(int(v) for v in np.unique(values) if int(v) > 0)


def choose_gallery_sources(by_source, limit=8):
    source_info = []
    for source_id, rows in by_source.items():
        labels = classes_present(Path(rows[0]["mask_path"]))
        source_info.append((source_id, rows, labels))
    selected = []
    covered = set()
    # Prioritize rare/informative classes then maximize remaining class diversity.
    for target in (6, 7, 8, 1, 2, 3, 4, 5):
        candidates = [item for item in source_info if target in item[2] and item[0] not in {v[0] for v in selected}]
        if candidates:
            candidates.sort(key=lambda item: (len(item[2] - covered), len(item[2])), reverse=True)
            selected.append(candidates[0]); covered |= candidates[0][2]
        if len(selected) == limit:
            return selected
    remaining = [item for item in source_info if item[0] not in {v[0] for v in selected}]
    remaining.sort(key=lambda item: (len(item[2] - covered), len(item[2])), reverse=True)
    for item in remaining:
        selected.append(item); covered |= item[2]
        if len(selected) == limit:
            break
    return selected


def class_list(labels):
    return ", ".join(NAMES[index] for index in sorted(labels))


def add_legend(slide, y):
    text(slide, "MINERAL LEGEND", .34, y, 1.18, .16, 9.5, NAVY, True)
    for index, (name, colour) in enumerate(zip(NAMES, PALETTE)):
        column = index % 5; row = index // 5
        x = 1.58 + column*2.30; yy = y+row*.22
        box(slide, x, yy, .13, .13, colour, rounded=False)
        text(slide, name, x+.18, yy-.01, 1.65, .15, 8.5, INK, True)


def stage(slide, title, body, x, y, w, colour):
    box(slide, x, y, w, 1.03, PALE, colour)
    box(slide, x, y, w, .27, colour, rounded=False)
    text(slide, title, x+.04, y+.05, w-.08, .13, 9.8, (255,255,255), True, PP_ALIGN.CENTER)
    text(slide, body, x+.11, y+.37, w-.22, .46, 12.5, INK, True, PP_ALIGN.CENTER)


def figure_one(prs, by_source, assets):
    preferred = by_source.get("Muscovite-Granite-3")
    source_id, records = ("Muscovite-Granite-3", preferred) if preferred else next(iter(by_source.items()))
    image_records = [records[0], records[len(records)//2], records[-1]]
    mask_png = assets / "fig1_shared_source_mask.png"; rgb_mask(Path(records[0]["mask_path"]), mask_png)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 1, "Dataset construction and source-level study design", "Original workflow using this study’s registered mineral images, indexed masks, and audited source-level manifest.")
    x_positions = [.30, 3.50, 6.70, 9.90]
    for i, (x, record) in enumerate(zip(x_positions, image_records), start=1):
        panel(slide, record["image_path"], f"REAL REGISTERED VIEW {frame_number(record['image_path'])}", x, 1.00, 2.85, 1.70, TEAL)
    panel(slide, mask_png, "ONE SHARED INDEXED MASK", x_positions[3], 1.00, 2.85, 1.70, RUST)
    text(slide, f"Example source: {source_id}. Eight angular views of the same thin section share registered geometry and one semantic reference mask.", .45, 3.08, 12.1, .22, 12.5, INK, True, PP_ALIGN.CENTER)
    stages = [
        ("1  ACQUIRE", "97 SOURCES\n776 RGB IMAGES", .35, BLUE),
        ("2  REGISTER", "8 VIEWS/SOURCE\nregistered geometry", 2.93, TEAL),
        ("3  ANNOTATE", "1 MASK/SOURCE\n8 mineral classes", 5.51, RUST),
        ("4  SPLIT SOURCES", "72 / 15 / 10\ntrain / val / test", 8.09, NAVY),
        ("5  VALIDATE", "2 × 5 FOLDS\n10 grouped estimates", 10.67, TEAL),
    ]
    for title, body, x, colour in stages:
        stage(slide, title, body, x, 3.75, 2.30, colour)
    for i in range(4):
        arrow(slide, 2.67+i*2.58, 4.27, 2.90+i*2.58, 4.27)
    text(slide, "Source-level rule: every registered view from one thin section is kept together in exactly one partition; correlated angular views cannot cross into testing.", .50, 5.16, 12.2, .23, 13.2, NAVY, True, PP_ALIGN.CENTER)
    box(slide, .75, 5.65, 11.8, .72, (231,242,244), TEAL)
    text(slide, "Primary evidence: validation-only model selection followed by repeated grouped source-level cross-validation", 1.03, 5.80, 11.2, .18, 14, INK, True, PP_ALIGN.CENTER)
    text(slide, "Dataset citation: Zhu and Yang (2026), Registered Angular Mineral Segmentation Dataset [1]. This study; high-level workflow organization is original.", .55, 6.93, 12.0, .15, 8.5, MID, False, PP_ALIGN.CENTER)


def figure_two(prs, by_source, assets):
    selected = choose_gallery_sources(by_source)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 2, "Representative registered mineral-image and semantic-mask gallery", "Eight real source examples selected to expose mineral, lithological, and optical diversity rather than to optimize appearance.")
    for i, (source_id, records, labels) in enumerate(selected):
        row, col = divmod(i, 4)
        x = .33 + col*3.25; y = .98 + row*2.77
        mask_png = assets / f"gallery_mask_{i+1:02d}.png"; rgb_mask(Path(records[0]["mask_path"]), mask_png)
        record = records[0]
        panel(slide, record["image_path"], "REAL POLARIZED-LIGHT IMAGE", x, y, 2.88, .93, TEAL, 7.8)
        panel(slide, mask_png, "INDEXED SEMANTIC MASK", x, y+1.28, 2.88, .93, RUST, 7.8)
        text(slide, source_id, x+.02, y+2.51, 2.84, .16, 10.2, INK, True, PP_ALIGN.CENTER)
        text(slide, f"{len(labels)} annotated mineral classes", x+.02, y+2.66, 2.84, .12, 8.6, MID, False, PP_ALIGN.CENTER)
    add_legend(slide, 6.57)
    text(slide, "Each pair retains the original source image and its real semantic mask. The gallery samples different sources; it is not a random tile collage or a set of model predictions.", .45, 7.03, 12.3, .15, 8.5, MID, False, PP_ALIGN.CENTER)


def figure_three(prs, by_source):
    preferred = by_source.get("Muscovite-Granite-3") or next(iter(by_source.values()))
    record = preferred[0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 3, "B7–B8 probability ensemble and validation-only selection", "Concise model framework: single-view inference, complementary architectures, probability fusion, and strictly separated model selection.")
    panel(slide, record["image_path"], "REAL INPUT IMAGE", .32, 1.10, 2.35, 1.62, TEAL)
    box(slide, .35, 3.04, 2.29, .46, PALE, TEAL)
    text(slide, "640 × 512 normalized input  |  448 × 448 training crop", .42, 3.12, 2.15, .22, 9.8, INK, True, PP_ALIGN.CENTER)
    stage(slide, "B7  PRETRAINED RESNET-34 U-NET", "Weighted cross-entropy\n+ mineral Dice loss", 3.12, 1.20, 2.72, BLUE)
    stage(slide, "B8  DEEPLABV3–RESNET50", "Weighted cross-entropy\n+ focal Tversky loss", 3.12, 3.62, 2.72, RUST)
    arrow(slide, 2.72, 2.00, 3.08, 1.72); arrow(slide, 2.72, 2.00, 3.08, 4.14)
    box(slide, 6.28, 1.84, 3.13, 1.58, (231,242,244), TEAL)
    text(slide, "PROBABILITY-LEVEL ENSEMBLE", 6.43, 2.08, 2.82, .16, 11, NAVY, True, PP_ALIGN.CENTER)
    text(slide, "pᵉⁿˢ = w pᴮ⁷ + (1 − w) pᴮ⁸\nŷ = argmax₍c₎ pᵉⁿˢ₍c₎", 6.54, 2.48, 2.60, .45, 16, INK, True, PP_ALIGN.CENTER)
    arrow(slide, 5.87, 1.72, 6.24, 2.38); arrow(slide, 5.87, 3.90, 6.24, 2.88)
    stage(slide, "VALIDATION-ONLY CALIBRATION", "Original split: B7 = 0.6\nB8 = 0.4", 9.84, 1.20, 2.90, NAVY)
    stage(slide, "GROUPED CV RESELECTION", "Fold B7 weights: 0.3–0.6\nmean = 0.44", 9.84, 3.38, 2.90, TEAL)
    arrow(slide, 9.45, 2.38, 9.80, 1.72); arrow(slide, 9.45, 2.88, 9.80, 3.90)
    box(slide, .66, 5.45, 12.0, .78, PALE, RUST)
    text(slide, "Model-selection safeguard: test sources never selected a checkpoint, crop policy, loss configuration, or ensemble coefficient.", .98, 5.68, 11.35, .20, 14.2, INK, True, PP_ALIGN.CENTER)
    text(slide, "Mathematical definitions, loss parameters, and grouped-CV uncertainty calculation are reported in Section 2.7; model architecture foundations are cited in [7,8].", .55, 6.89, 12.2, .15, 8.5, MID, False, PP_ALIGN.CENTER)


def add_bar(slide, label, value, x, y, width, maximum, colour):
    text(slide, label, x, y-.015, 1.46, .20, 9.2, INK, True)
    box(slide, x+1.50, y+.02, width, .14, (222,232,236), rounded=False)
    box(slide, x+1.50, y+.02, width*value/maximum, .14, colour, rounded=False)
    text(slide, f"{value:.4f}", x+1.50+width+.08, y-.02, .48, .17, 8.7, INK, True)


def figure_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 4, "Model-development progression and grouped source-level validation", "All bars, uncertainty annotations, class labels, and values are editable PowerPoint objects. Values are validated study results.")
    text(slide, "A. Original fixed source-held test: mean mineral IoU", .40, 1.00, 5.0, .20, 14, NAVY, True)
    progression = [("B0 compact U-Net", .0837, BLUE), ("B1 class weights", .0908, BLUE), ("B2 mineral crops", .1296, TEAL), ("B7 ResNet-34", .2138, NAVY), ("B8 DeepLabV3", .2096, RUST), ("B7–B8 ensemble", .2440, TEAL)]
    for i, (label, value, colour) in enumerate(progression):
        add_bar(slide, label, value, .48, 1.43+i*.43, 2.35, .27, colour)
    box(slide, 5.12, 1.02, 3.28, 2.98, (231,242,244), TEAL)
    text(slide, "B. PRIMARY GROUPED-CV ESTIMATE", 5.28, 1.23, 2.96, .18, 11.2, NAVY, True, PP_ALIGN.CENTER)
    text(slide, "2 repeats × 5 source-level folds", 5.31, 1.58, 2.90, .18, 12.3, INK, True, PP_ALIGN.CENTER)
    text(slide, "Mineral mIoU\n0.3294 ± 0.0459", 5.26, 2.05, 1.45, .54, 15, INK, True, PP_ALIGN.CENTER)
    text(slide, "Mineral Dice\n0.4738 ± 0.0545", 6.83, 2.05, 1.45, .54, 15, INK, True, PP_ALIGN.CENTER)
    text(slide, "Approx. 95% mIoU: 0.3009–0.3578\nApprox. 95% Dice: 0.4401–0.5076", 5.31, 3.28, 2.90, .40, 11, MID, False, PP_ALIGN.CENTER)
    text(slide, "C. Per-mineral mean IoU across grouped held-out folds", 8.88, 1.00, 4.0, .20, 14, NAVY, True)
    class_results = [("Olivine", .3701), ("Pyroxene", .2413), ("Plagioclase", .2974), ("Alkali feldspar", .4035), ("Quartz", .3482), ("Biotite", .1796), ("Muscovite", .4142), ("Hornblende", .3807)]
    for i, (label, value) in enumerate(class_results):
        column, row = divmod(i, 4)
        x = 8.88 + column*2.12; y = 1.43 + row*.51
        colour = TEAL if value >= .30 else RUST
        box(slide, x, y+.03, .14, .14, colour, rounded=False)
        text(slide, label, x+.20, y-.015, 1.23, .18, 9.0, INK, True)
        text(slide, f"{value:.3f}", x+1.46, y-.015, .48, .18, 9.0, INK, True, PP_ALIGN.RIGHT)
    box(slide, .55, 4.65, 12.15, 1.18, PALE, NAVY)
    text(slide, "Finding", .82, 4.87, 1.02, .18, 12.2, NAVY, True)
    text(slide, "The ensemble improved the original source-held test relative to B7 and B8. Cross-validation is the primary performance estimate because it summarizes ten independent source-level hold-out evaluations. Per-class variation is retained rather than concealed by a single macro score.", 1.74, 4.75, 10.56, .52, 12.5, INK, False)
    text(slide, "Full confusion matrix, class distribution, normalized palette, and fold-specific weights are reported in Supplementary Figures S1–S4.", .55, 6.78, 12.0, .16, 8.7, MID, False, PP_ALIGN.CENTER)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--output", default="artifacts/mineral_segmentation_b7_b8_ensemble/expanded_context_figures")
    args = parser.parse_args()
    with Path(args.manifest).open(newline="", encoding="utf-8") as handle:
        rows = list(csv.DictReader(handle))
    if not rows:
        raise ValueError("Manifest is empty.")
    for required in ("source_id", "image_path", "mask_path"):
        if required not in rows[0]:
            raise KeyError(f"Manifest does not contain required column: {required}")
    by_source = source_records(rows)
    output = Path(args.output); output.mkdir(parents=True, exist_ok=True)
    assets = output / "context_figure_real_assets"; assets.mkdir(parents=True, exist_ok=True)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    figure_one(prs, by_source, assets)
    figure_two(prs, by_source, assets)
    figure_three(prs, by_source)
    figure_four(prs)
    deck = output / "Expanded_Mineral_Context_Figures_EDITABLE.pptx"
    prs.save(deck)
    (output / "expanded_context_figures_summary.txt").write_text(
        "Created four original editable context figures from the study manifest.\n"
        "Figure 1: dataset construction and source-level study design.\n"
        "Figure 2: representative real-image and semantic-mask gallery.\n"
        "Figure 3: B7-B8 framework and validation-only selection logic.\n"
        "Figure 4: model development and grouped validation evidence.\n",
        encoding="utf-8")
    print(deck)


if __name__ == "__main__":
    main()
