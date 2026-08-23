"""Build a truly editable PowerPoint master for the B7-B8 ensemble figures.

Only the scientific microscope images, masks and predictions are raster. Every title,
panel heading, source label, IoU callout, legend entry, workflow box, arrow and chart
annotation is created as a separate PowerPoint shape or text object.
"""
from __future__ import annotations
import argparse,csv,re,json
from pathlib import Path
import numpy as np
import torch
from PIL import Image
from torchvision.models.segmentation import deeplabv3_resnet50
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE,MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from generate_b7_native_900dpi_figures import PretrainedResUNet,MEAN,STD,NAMES,PALETTE,miou,rgb
NAVY=(18,73,110);TEAL=(26,126,111);BLUE=(50,120,168);RUST=(180,92,72);INK=(20,34,45);PALE=(244,248,250)
def num(path):
 m=re.search(r'_(\d+)\.(?:jpg|jpeg|png)$',path.lower());return int(m.group(1)) if m else 0
def rect(s,x,y,w,h,fill,line=None,round=True):
 sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=RGBColor(*fill);sh.line.color.rgb=RGBColor(*(line or fill));return sh
def txt(s,v,x,y,w,h,sz=14,c=INK,b=False,a=PP_ALIGN.LEFT):
 sh=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE;p=tf.paragraphs[0];p.alignment=a;r=p.add_run();r.text=v;r.font.name='Aptos';r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=RGBColor(*c);return sh
def arrow(s,x1,y1,x2,y2):
 l=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2));l.line.color.rgb=RGBColor(*TEAL);l.line.width=Pt(2.6);l.line.end_arrowhead=True
def header(s,title,subtitle=''):
 rect(s,0,0,13.333,.55,NAVY,round=False);txt(s,title,.28,.07,12.5,.20,21,(255,255,255),True);txt(s,subtitle,.30,.60,12.3,.18,10,(75,91,102))
def ensemble(b7,b8,im,d,w):
 x=torch.from_numpy(np.asarray(im,dtype=np.float32).transpose(2,0,1)/255.).unsqueeze(0);x=(x-MEAN)/STD
 with torch.no_grad():return (w*b7(x.to(d)).softmax(1)+(1-w)*b8(x.to(d))['out'].softmax(1)).argmax(1).squeeze().cpu().numpy().astype(np.int64)
def add_legend(s,y):
 txt(s,'MINERAL LEGEND',.30,y,1.1,.16,10,NAVY,True)
 for i,(name,col) in enumerate(zip(NAMES,PALETTE)):
  row=i//5;pos=i%5;x=1.55+pos*2.30;yy=y+row*.25
  rect(s,x,yy,.14,.14,col,round=False);txt(s,name,x+.20,yy-.01,1.75,.16,9,INK,True)
def panel(s,img,label,x,y,w,h,col):
 rect(s,x,y,w,.30,col,round=False);txt(s,label,x+.02,y+.07,w-.04,.14,10,(255,255,255),True,PP_ALIGN.CENTER);s.shapes.add_picture(str(img),Inches(x),Inches(y+.34),width=Inches(w),height=Inches(h))
def main():
 p=argparse.ArgumentParser();p.add_argument('--manifest',required=True);p.add_argument('--b7-checkpoint',required=True);p.add_argument('--b8-checkpoint',required=True);p.add_argument('--b7-weight',type=float,default=.6);p.add_argument('--output',required=True);a=p.parse_args()
 if not torch.cuda.is_available():raise RuntimeError('CUDA required')
 out=Path(a.output);out.mkdir(parents=True,exist_ok=True);assets=out/'editable_ppt_native_assets';assets.mkdir(parents=True,exist_ok=True)
 with Path(a.manifest).open(newline='',encoding='utf-8') as f:rows=[r for r in csv.DictReader(f) if r['split']=='test']
 d=torch.device('cuda');b7=PretrainedResUNet().to(d);b8=deeplabv3_resnet50(weights=None,weights_backbone=None,num_classes=9,aux_loss=False).to(d)
 b7.load_state_dict(torch.load(a.b7_checkpoint,map_location=d,weights_only=False)['model']);b8.load_state_dict(torch.load(a.b8_checkpoint,map_location=d,weights_only=False)['model']);b7.eval();b8.eval();records=[]
 for i,r in enumerate(rows):
  with Image.open(r['image_path']) as x:im=x.convert('RGB')
  with Image.open(r['mask_path']) as x:gt=np.asarray(x.convert('L'),dtype=np.int64)
  pr=ensemble(b7,b8,im,d,a.b7_weight);base=assets/f'case_{i:03d}';im.save(base.with_name(base.name+'_image.png'));rgb(gt).save(base.with_name(base.name+'_gt.png'));rgb(pr).save(base.with_name(base.name+'_pred.png'))
  records.append({'source_id':r['source_id'],'image_path':r['image_path'],'image_file':base.with_name(base.name+'_image.png'),'gt_file':base.with_name(base.name+'_gt.png'),'pred_file':base.with_name(base.name+'_pred.png'),'gt':gt,'miou':miou(gt,pr)})
 records.sort(key=lambda z:z['miou']);groups={'Strong':records[-2:],'Typical':records[len(records)//2:len(records)//2+2],'Difficult':records[:2]};best=groups['Strong'][-1]
 same=sorted([r for r in rows if r['source_id']==best['source_id']],key=lambda r:num(r['image_path']))
 for tag,r in [('view1',same[0]),('view8',same[-1])]:
  with Image.open(r['image_path']) as im:im.convert('RGB').save(assets/f'{tag}.png')
 # Build deck.
 prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)
 # Slide 1 data.
 s=prs.slides.add_slide(prs.slide_layouts[6]);header(s,'Figure 1. Registered mineral data','All labels, source names and legend items are editable. Only scientific image pixels are raster data.')
 for x,lab,path,col in [(0.35,'REGISTERED VIEW 1',assets/'view1.png',TEAL),(4.55,'REGISTERED VIEW 8',assets/'view8.png',TEAL),(8.75,'INDEXED SOURCE MASK',best['gt_file'],RUST)]:panel(s,path,lab,x,1.05,3.75,3.0,col)
 txt(s,f"Source: {best['source_id']}  |  Eight registered angular views per source  |  Native data: 1280 × 1024 pixels",.38,4.48,12.2,.23,13,INK,True);add_legend(s,5.05)
 # Slide 2 workflow.
 s=prs.slides.add_slide(prs.slide_layouts[6]);header(s,'Figure 2. Leakage-safe ensemble workflow','Sparse methodology figure: each workflow component is an editable PowerPoint object.')
 stages=[('1  GROUP SOURCES','8 registered views\n1 indexed mask',BLUE,.45,1.15),('2  SPLIT SOURCES','72 train | 15 validation\n10 held-out test',TEAL,3.35,1.15),('3  TRAIN B7 + B8','Pretrained ResNet-34 U-Net\nDeepLabV3-ResNet50',RUST,6.25,1.15),('4  CALIBRATE','Coefficient selected\non validation only',NAVY,2.00,4.10),('5  GROUPED CV','2 repeats × 5 folds\nmIoU = 0.3294 ± 0.0459',TEAL,5.70,4.10)]
 for title,body,col,x,y in stages:rect(s,x,y,2.55,1.35,PALE,line=col);rect(s,x,y,2.55,.34,col,round=False);txt(s,title,x+.04,y+.08,2.47,.15,12,(255,255,255),True,PP_ALIGN.CENTER);txt(s,body,x+.14,y+.53,2.25,.54,14,INK,True,PP_ALIGN.CENTER)
 arrow(s,3.0,1.83,3.32,1.83);arrow(s,5.9,1.83,6.22,1.83);arrow(s,7.52,2.55,7.00,3.94);arrow(s,4.56,4.78,5.67,4.78)
 txt(s,'Source-level grouping prevents views of one thin section from appearing in more than one partition.',.50,6.35,12.1,.25,15,NAVY,True)
 # Slide 3 validation.
 s=prs.slides.add_slide(prs.slide_layouts[6]);header(s,'Figure 3. Repeated source-level validation of the B7-B8 ensemble','Ten outer held-out estimates; ensemble weight selected separately inside every fold using validation sources only.')
 cards=[('MINERAL mIoU','0.3294 ± 0.0459','Approx. 95% CI: 0.3009–0.3578',NAVY,.45),('MINERAL DICE','0.4738 ± 0.0545','Approx. 95% CI: 0.4401–0.5076',TEAL,4.52),('MODEL SELECTION','B7 + B8 ensemble','Mean B7 probability weight: 0.44',RUST,8.59)]
 for title,big,sub,col,x in cards:rect(s,x,1.35,3.75,2.10,PALE,line=col);rect(s,x,1.35,3.75,.38,col,round=False);txt(s,title,x+.03,1.44,3.68,.18,13,(255,255,255),True,PP_ALIGN.CENTER);txt(s,big,x+.1,2.15,3.55,.45,24,INK,True,PP_ALIGN.CENTER);txt(s,sub,x+.1,2.88,3.55,.24,13,INK,True,PP_ALIGN.CENTER)
 txt(s,'Per-mineral grouped-fold mean IoU',.55,4.12,3.4,.22,16,NAVY,True)
 vals=[('Olivine',.3701),('Pyroxene',.2413),('Plagioclase',.2974),('Alkali feldspar',.4035),('Quartz',.3482),('Biotite',.1796),('Muscovite',.4142),('Hornblende',.3807)]
 for i,(n,v) in enumerate(vals):
  x=.65+(i%4)*3.15;y=4.60+(i//4)*.65;txt(s,n,x,y,1.35,.18,11,INK,True);rect(s,x+1.38,y+.03,1.35,.14,(222,232,236),round=False);rect(s,x+1.38,y+.03,1.35*v/.45,.14,TEAL,round=False);txt(s,f'{v:.3f}',x+2.80,y,0.5,.18,10,INK,True)
 # qualitative slides
 for group,recs in groups.items():
  s=prs.slides.add_slide(prs.slide_layouts[6]);subtitle={'Strong':'Favourable held-out cases. Visual agreement is paired with the measured frame mIoU.','Typical':'Central-ranked held-out cases. Broad domains are recovered with remaining class and boundary confusion.','Difficult':'Low-ranked held-out cases. These are retained to expose composition-dependent limitations.'}[group];header(s,f'Figure {4 if group=="Strong" else 5 if group=="Typical" else 6}. {group} ensemble predictions',subtitle)
  for rix,r in enumerate(recs):
   y=1.05+rix*3.02;col=TEAL if rix==0 else RUST
   panel(s,r['image_file'],'POLARIZED-LIGHT IMAGE',.32,y,3.75,2.38,col);panel(s,r['gt_file'],'INDEXED GROUND TRUTH',4.78,y,3.75,2.38,RUST);panel(s,r['pred_file'],'B7-B8 ENSEMBLE PREDICTION',9.24,y,3.75,2.38,RUST)
   txt(s,f"{chr(65+rix)}  Source: {r['source_id']}  |  ensemble frame mineral mIoU = {r['miou']:.3f}",.38,y+2.72,12.2,.18,13,INK,True)
  add_legend(s,7.18)
 ppt=out/'B7_B8_Ensemble_TRULY_EDITABLE_Master.pptx';prs.save(ppt)
 summary={'deck':str(ppt),'editable_elements':'all titles, panel labels, source labels, metric callouts, legends, workflow boxes, arrows, validation cards and bars','raster_elements':'native microscope images, semantic masks and ensemble predictions only','b7_weight':a.b7_weight,'b8_weight':1-a.b7_weight,'slides':6};(out/'truly_editable_pptx_summary.json').write_text(json.dumps(summary,indent=2),encoding='utf-8');print(json.dumps(summary,indent=2))
if __name__=='__main__':main()
